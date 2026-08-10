"""Generate OnScen presentation (non-technical, French) — sponsors, lieux, créateurs."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent / "OnScen-Presentation.pptx"
LEGACY_OUTPUT = Path(__file__).resolve().parent / "Soundly-Presentation.pptx"
SCREENSHOTS = Path(__file__).resolve().parent.parent / "presentation-screenshots"

# Dark music-app palette (aligné UI OnScen)
BG = RGBColor(0x0D, 0x0D, 0x14)
BG_ALT = RGBColor(0x14, 0x14, 0x22)
ACCENT = RGBColor(0xA7, 0x8B, 0xFA)  # violet
ACCENT2 = RGBColor(0xF4, 0x72, 0xB6)  # pink
TEXT = RGBColor(0xF4, 0xF4, 0xF8)
MUTED = RGBColor(0xA8, 0xA8, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDES = [
    {
        "type": "title",
        "title": "OnScen",
        "subtitle": "Réseau social musique live & sorties",
        "tagline": "Découvrir, écouter ensemble et sortir — sur une carte, en salon ou en live · getsoundy.com",
        "image": "01-actualite.png",
        "image_caption": "Actualités, stories, reels — navigation en 5 onglets",
    },
    {
        "title": "Le problème",
        "bullets": [
            "Musique dispersée entre TikTok, Spotify, Instagram et Shotgun — aucun parcours unifié",
            "Difficile de voir qui sort où ce soir, près de chez soi",
            "Bars et salles 50–500 places : peu de visibilité entre deux dates",
            "Créateurs mid-tier sans hub simple entre live, communauté et monétisation",
        ],
        "footer": "OnScen reconnecte écoute sociale, carte et événements",
    },
    {
        "title": "La solution",
        "highlight": "OnScen est le réseau social qui fait écouter la même musique au même moment — sur une carte géolocalisée, en salon d'écoute ou en live vidéo.",
        "image": "02-carte.png",
        "image_layout": "right",
        "image_caption": "Carte : événements, salons, lives et filtres Autour / Pays",
    },
    {
        "title": "Comment ça marche ?",
        "steps": [
            ("1", "Profil & affinités", "Goûts musicaux, ville, abonnements créateurs"),
            ("2", "Carte & agenda", "Événements 3 jours, salons et lives à proximité"),
            ("3", "Rejoindre ou lancer", "Salon YouTube sync ou live caméra + chat"),
            ("4", "Partager & soutenir", "Stories, reels, pourboires Stripe Connect"),
        ],
        "image": "07-messages.png",
        "image_caption": "Messages, matchs et groupes",
    },
    {
        "title": "Carte & événements",
        "bullets": [
            "Pins événements : aujourd'hui, demain et 3 prochains jours (filtre Événement)",
            "Section ✨ Sponso : carrousel partenaires en tête de sidebar",
            "Salons d'écoute et lives géolocalisés · vue globe 3D satellite",
            "Onglets Autour (proximité) et Pays (agenda national)",
            "Respect vie privée : position floutée · contrôle de visibilité",
        ],
        "image": "02-carte.png",
        "image_caption": "Carte interactive · panneau browse événements",
    },
    {
        "title": "Salons & petits salons",
        "bullets": [
            "Watch party musicale synchronisée (YouTube)",
            "Chat, file d'attente, modération hôte",
            "Petit salon : rejoindre depuis la carte sans quitter la vue",
            "Salons publics ou privés · ancrés sur la carte",
            "Idéal soirées entre amis ou communautés locales",
        ],
        "image": "04-salon.png",
        "image_caption": "Grand salon — lecteur, file d'attente, propositions",
    },
    {
        "title": "Lives musicaux",
        "bullets": [
            "Diffusion caméra + musique en direct (LiveKit / Cloudflare HLS)",
            "Chat temps réel, réactions et pourboires visuels",
            "Pourboires & abonnements créateurs via Stripe Connect (prod)",
            "Grille découverte · intégration théâtre / PiP",
            "16+ pour performer · 18+ pour encaisser des pourboires",
        ],
        "image": "06-live.png",
        "image_secondary": "05-lives-tab.png",
        "image_caption": "Live avec chat · grille Direct",
    },
    {
        "title": "Pour les lieux & sponsors",
        "bullets": [
            "7 emplacements natifs : carte Sponso, bandeau, fil, stories, reels, théâtre salon",
            "Idéal bars, clubs, festivals : pin carte + carrousel événement sponsorisé",
            "Tarifs transparents dès ~150 € HT / 7 jours (palier lancement)",
            "Offre pilote lieux : Sidebar Sponso 7 j · 99 € HT (10 places)",
            "Estimation d'audience avant campagne · badge Sponsorisé visible",
        ],
        "image": "02-carte.png",
        "image_caption": "Visibilité locale : pin événement + section Sponso",
    },
    {
        "title": "Pour les créateurs",
        "bullets": [
            "Profil public : reels, abonnés, favoris, événements",
            "Live + pourboires + abonnements Supporter / Super fan",
            "Stripe Connect opérationnel · OnScen+ (soutien plateforme)",
            "13+ compte · 16+ live · 18+ monétisation",
            "Plateforme sociale musicale — pas un service de rencontre",
        ],
        "image": "08-profil.png",
        "image_caption": "Profil créateur · reels et communauté",
    },
    {
        "title": "Sécurité & conformité",
        "bullets": [
            "Âge minimum 13 ans · autorisation parentale 13–18",
            "Paiements Stripe — aucune carte stockée chez OnScen",
            "RGPD : export données, suppression compte, privacy@getsoundy.com",
            "Modération : signalement, bannissement, CGU · hébergement France (Scaleway)",
            "Formats sponsorisés conformes DSA (transparence publicitaire)",
        ],
    },
    {
        "title": "Disponible partout",
        "bullets": [
            "Web & PWA : getsoundy.com · installation écran d'accueil",
            "Applications natives iOS & Android (Capacitor)",
            "Une expérience mobile-first · temps réel (chat, sync salon)",
            "Production opérationnelle — juillet 2026",
        ],
        "image": "01-actualite.png",
        "image_caption": "Même expérience web, PWA et mobile",
    },
    {
        "title": "Vision",
        "bullets": [
            "Densifier la communauté musique ville par ville (France d'abord)",
            "Accélérer partenariats lieux & sponsors fondateurs",
            "Renforcer découverte locale : carte, événements, reels",
            "OnScen : où la musique devient un moment partagé — en ligne et sur place",
        ],
        "image": "03-petit-salon.png",
        "image_caption": "Petit salon sur la carte — écoute rapide",
    },
    {
        "type": "contact",
        "title": "Contact & prochaine étape",
        "bullets": [
            "Tester OnScen : getsoundy.com (web ou app)",
            "Organisateur / bar : demander devis Sponso · contact@getsoundy.com",
            "Créateur : lancer un salon ou un live près de chez vous",
            "Partenariat : one-pager commercial & estimation audience sur demande",
        ],
        "cta": "Prêts à écouter — et sortir — ensemble ?",
        "image": "05-lives-tab.png",
        "image_caption": "Rejoignez un live ou publiez votre prochaine date",
    },
]


def img_path(name: str | None) -> Path | None:
    if not name:
        return None
    path = SCREENSHOTS / name
    return path if path.is_file() else None


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, left=Inches(0), top=Inches(0), width=Inches(0.12), height=Inches(7.5)):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_gradient_bar(slide):
    bar = slide.shapes.add_shape(1, Inches(0.6), Inches(1.35), Inches(8.8), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT2
    bar.line.fill.background()


def style_paragraph(p, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.alignment = align
    p.space_after = Pt(8)


def add_title_block(slide, title, top=Inches(0.55)):
    box = slide.shapes.add_textbox(Inches(0.75), top, Inches(8.5), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    style_paragraph(p, size=32, bold=True, color=WHITE)
    add_gradient_bar(slide)
    return box


def add_bullets(slide, items, left=Inches(0.85), top=Inches(1.75), width=Inches(8.3), size=20):
    box = slide.shapes.add_textbox(left, top, width, Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        style_paragraph(p, size=size, color=TEXT)
        p.space_before = Pt(6)
    return box


def add_image_frame(slide, path: Path, left, top, width, caption: str | None = None):
    """Add screenshot with subtle phone-frame border."""
    frame = slide.shapes.add_shape(
        1,
        left - Inches(0.06),
        top - Inches(0.06),
        width + Inches(0.12),
        width * 1.95 + Inches(0.12),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = BG_ALT
    frame.line.color.rgb = ACCENT
    frame.line.width = Pt(1.5)

    pic_h = width * 1.85
    slide.shapes.add_picture(str(path), left, top, width=width, height=pic_h)

    if caption:
        cap = slide.shapes.add_textbox(left, top + pic_h + Inches(0.08), width + Inches(0.2), Inches(0.55))
        p = cap.text_frame.paragraphs[0]
        p.text = caption
        style_paragraph(p, size=11, color=MUTED, align=PP_ALIGN.CENTER)


def add_slide_image(slide, data, default_left=Inches(0.55), default_top=Inches(1.55), width=Inches(2.65)):
    path = img_path(data.get("image"))
    if not path:
        return False

    layout = data.get("image_layout", "left")
    if layout == "right":
        left = Inches(6.55)
    else:
        left = default_left

    add_image_frame(slide, path, left, default_top, width, data.get("image_caption"))

    secondary = img_path(data.get("image_secondary"))
    if secondary:
        add_image_frame(
            slide,
            secondary,
            left,
            default_top + Inches(3.35),
            width * 0.92,
            None,
        )
    return True


def build_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)

    circle = slide.shapes.add_shape(9, Inches(6.8), Inches(0.4), Inches(2.2), Inches(2.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BG_ALT
    circle.line.color.rgb = ACCENT
    circle.line.width = Pt(2)

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(5.8), Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    p.text = data["title"]
    style_paragraph(p, size=54, bold=True, color=WHITE)

    sub_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.65), Inches(5.8), Inches(0.8))
    p = sub_box.text_frame.paragraphs[0]
    p.text = data["subtitle"]
    style_paragraph(p, size=26, bold=False, color=ACCENT)

    tag_box = slide.shapes.add_textbox(Inches(0.9), Inches(3.55), Inches(5.5), Inches(1.0))
    p = tag_box.text_frame.paragraphs[0]
    p.text = data.get("tagline", "")
    style_paragraph(p, size=17, color=MUTED)

    path = img_path(data.get("image"))
    if path:
        add_image_frame(slide, path, Inches(6.35), Inches(1.35), Inches(2.85), data.get("image_caption"))


def build_highlight_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_block(slide, data["title"])

    has_image = add_slide_image(slide, data, default_top=Inches(1.7), width=Inches(2.5))
    text_left = Inches(0.9)
    text_width = Inches(8.2)
    if has_image and data.get("image_layout") == "right":
        text_width = Inches(5.4)
    elif has_image:
        text_left = Inches(3.55)
        text_width = Inches(5.8)

    box = slide.shapes.add_textbox(text_left, Inches(2.0), text_width, Inches(3.8))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = data["highlight"]
    style_paragraph(p, size=24, color=WHITE, align=PP_ALIGN.LEFT)
    p.line_spacing = 1.35


def build_steps_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_ALT)
    add_accent_bar(slide)
    add_title_block(slide, data["title"])

    has_image = add_slide_image(slide, data, default_left=Inches(6.2), default_top=Inches(1.55), width=Inches(2.55))
    step_width = Inches(5.2) if has_image else Inches(7.5)

    y_start = 1.85
    for i, (num, title, desc) in enumerate(data["steps"]):
        y = y_start + i * 1.15
        badge = slide.shapes.add_shape(9, Inches(0.85), Inches(y), Inches(0.55), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()
        nb = badge.text_frame
        nb.paragraphs[0].text = num
        style_paragraph(nb.paragraphs[0], size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        tbox = slide.shapes.add_textbox(Inches(1.6), Inches(y - 0.05), step_width, Inches(0.45))
        p = tbox.text_frame.paragraphs[0]
        p.text = title
        style_paragraph(p, size=21, bold=True, color=WHITE)

        dbox = slide.shapes.add_textbox(Inches(1.6), Inches(y + 0.35), step_width, Inches(0.4))
        p = dbox.text_frame.paragraphs[0]
        p.text = desc
        style_paragraph(p, size=15, color=MUTED)


def build_bullet_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_block(slide, data["title"])

    has_image = add_slide_image(slide, data)
    bullets_left = Inches(3.55) if has_image else Inches(0.85)
    bullets_width = Inches(5.9) if has_image else Inches(8.3)
    bullets = [f"•  {b}" for b in data.get("bullets", [])]
    add_bullets(slide, bullets, left=bullets_left, top=Inches(1.85), width=bullets_width)

    if data.get("footer"):
        fbox = slide.shapes.add_textbox(Inches(0.85), Inches(6.5), Inches(8.3), Inches(0.6))
        p = fbox.text_frame.paragraphs[0]
        p.text = data["footer"]
        style_paragraph(p, size=16, bold=True, color=ACCENT2)


def build_contact_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_ALT)
    add_accent_bar(slide)
    add_title_block(slide, data["title"])

    has_image = add_slide_image(slide, data, default_left=Inches(6.15), default_top=Inches(1.7), width=Inches(2.6))
    bullets_left = Inches(0.85)
    bullets_width = Inches(5.5) if has_image else Inches(8.3)

    bullets = [f"→  {b}" for b in data.get("bullets", [])]
    add_bullets(slide, bullets, left=bullets_left, top=Inches(1.9), width=bullets_width, size=21)

    cta = slide.shapes.add_textbox(Inches(0.85), Inches(5.8), bullets_width, Inches(0.8))
    p = cta.text_frame.paragraphs[0]
    p.text = data.get("cta", "")
    style_paragraph(p, size=24, bold=True, color=ACCENT, align=PP_ALIGN.LEFT if has_image else PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    missing = []
    for data in SLIDES:
        for key in ("image", "image_secondary"):
            name = data.get(key)
            if name and not img_path(name):
                missing.append(name)

    if missing:
        print(f"Warning: missing screenshots: {', '.join(sorted(set(missing)))}")

    for data in SLIDES:
        t = data.get("type", "bullet")
        if t == "title":
            build_title_slide(prs, data)
        elif "highlight" in data:
            build_highlight_slide(prs, data)
        elif "steps" in data:
            build_steps_slide(prs, data)
        elif t == "contact":
            build_contact_slide(prs, data)
        else:
            build_bullet_slide(prs, data)

    prs.save(str(OUTPUT))
    # Alias legacy filename for liens / scripts existants
    prs.save(str(LEGACY_OUTPUT))
    print(f"Created: {OUTPUT}")
    print(f"Created: {LEGACY_OUTPUT} (alias legacy)")
    print(f"Slides: {len(SLIDES)}")
    print(f"Screenshots dir: {SCREENSHOTS}")


if __name__ == "__main__":
    main()
