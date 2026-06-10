"""Generate Soundly presentation (non-technical, French)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent / "Soundly-Presentation.pptx"
SCREENSHOTS = Path(__file__).resolve().parent / "presentation-screenshots"

# Dark music-app palette
BG = RGBColor(0x0D, 0x0D, 0x14)
BG_ALT = RGBColor(0x14, 0x14, 0x22)
ACCENT = RGBColor(0xA7, 0x8B, 0xFA)  # violet
ACCENT2 = RGBColor(0xF4, 0x72, 0xB6)  # pink
TEXT = RGBColor(0xF4, 0xF4, 0xF8)
MUTED = RGBColor(0xA8, 0xA8, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDES = [
    {
        "type": "title",
        "title": "Soundly",
        "subtitle": "Écoute collective & lives musicaux",
        "tagline": "Découvrir, écouter et partager la musique — ensemble, près de chez vous",
        "image": "01-actualite.png",
        "image_caption": "Fil d’actualité, stories et navigation par onglets",
    },
    {
        "title": "Le problème",
        "bullets": [
            "Écouter seul sur Spotify ou YouTube, c’est bien… mais isolant",
            "Difficile de trouver qui partage vos goûts musicaux à proximité",
            "Les créateurs manquent d’espaces simples pour performer en live",
            "Pas de lien naturel entre carte, communauté et monétisation",
        ],
        "footer": "Une opportunité : reconnecter la musique au social et au local",
    },
    {
        "title": "La solution",
        "highlight": "Soundly est l’app qui vous fait écouter la même musique, au même moment, avec des gens autour de vous — sur une carte, en salon ou en live.",
        "image": "02-carte.png",
        "image_layout": "right",
        "image_caption": "Carte géolocalisée : salons, lives et personnes à proximité",
    },
    {
        "title": "Comment ça marche ?",
        "steps": [
            ("1", "Créez votre profil", "Goûts musicaux, photo, affinités"),
            ("2", "Explorez la carte", "Salons et lives près de vous"),
            ("3", "Rejoignez ou lancez", "Écoute synchronisée Spotify / YouTube"),
            ("4", "Partagez l’instant", "Chat, réactions, pourboires en live"),
        ],
        "image": "07-messages.png",
        "image_caption": "Messages privés, matchs et groupes",
    },
    {
        "title": "Carte & proximité",
        "bullets": [
            "Visualisez qui écoute quoi autour de vous",
            "Salons d’écoute et lives géolocalisés",
            "Position floutée (~50 m) pour protéger votre vie privée",
            "Filtres : distance, affinités musicales, lives uniquement",
            "Découvrez des profils et des communautés locales",
        ],
        "image": "02-carte.png",
        "image_caption": "Carte interactive avec panneau « À proximité »",
    },
    {
        "title": "Salons & petits salons",
        "bullets": [
            "Salons : watch party musicale à plusieurs",
            "Musique synchronisée via Spotify ou YouTube",
            "Chat en direct, playlist partagée, modération par l’hôte",
            "Petit salon : rejoindre rapidement depuis la carte",
            "Idéal pour soirées, amis ou rencontres autour d’un morceau",
        ],
        "image": "04-salon.png",
        "image_caption": "Grand salon — lecteur, file d’attente et propositions de morceaux",
    },
    {
        "title": "Lives musicaux",
        "bullets": [
            "Diffusez musique + caméra en direct (style Twitch / TikTok)",
            "Spectateurs en temps réel, chat et réactions gratuites",
            "Pourboires visuels style cadeaux live (1 € à 100 €)",
            "Caméra relayée en direct (WebRTC) — fluide et interactive",
            "Grille de découverte inspirée des plateformes de streaming",
        ],
        "image": "06-live.png",
        "image_secondary": "05-lives-tab.png",
        "image_caption": "Page live avec caméra, chat et pourboires · Grille de découverte",
    },
    {
        "title": "Pour les créateurs",
        "bullets": [
            "13+ : compte et participation aux salons",
            "16+ : lancer un live (performance artistique)",
            "18+ : recevoir pourboires et abonnements",
            "Profil public, reels, abonnés et statuts Supporter",
            "Soundly = plateforme sociale musicale, pas un service de rencontre",
        ],
        "image": "08-profil.png",
        "image_caption": "Profil public : photos, reels, abonnés et favoris",
    },
    {
        "title": "Monétisation",
        "bullets": [
            "Pourboires volontaires pendant les lives",
            "Abonnements créateurs : Supporter, Super fan…",
            "Soundly+ : soutien mensuel à la plateforme",
            "Paiements sécurisés via Stripe en production",
            "Simulation gratuite en version démo (aucune carte demandée)",
        ],
        "image": "09-monetisation.png",
        "image_caption": "Soundly+ et paramètres de compte",
    },
    {
        "title": "Sécurité & conformité",
        "bullets": [
            "Âge minimum : 13 ans (autorisation parentale 13–18)",
            "Stripe : saisie carte côté Stripe, jamais stockée chez Soundly",
            "RGPD : politique de confidentialité, suppression de compte",
            "Géolocalisation : contrôle de la visibilité (masquer, ville seule)",
            "Modération : signalement, bannissement, CGU claires",
        ],
    },
    {
        "title": "Disponible partout",
        "bullets": [
            "Web & PWA : navigateur, installation sur écran d’accueil",
            "iOS : application native (Capacitor)",
            "Android : application native (Capacitor)",
            "Une seule expérience, tous vos appareils",
            "Temps réel : synchronisation musique et chat instantané",
        ],
        "image": "01-actualite.png",
        "image_caption": "Même interface sur web, PWA et mobile natif",
    },
    {
        "title": "Vision & prochaines étapes",
        "bullets": [
            "Renforcer la communauté locale autour de la musique",
            "Ouvrir la monétisation créateurs en production",
            "Élargir les intégrations (playlists, découvertes)",
            "Finaliser conformité légale et lancement public",
            "Soundly : où la musique devient un moment partagé",
        ],
        "image": "03-petit-salon.png",
        "image_caption": "Petit salon sur la carte — écoute rapide sans quitter la carte",
    },
    {
        "type": "contact",
        "title": "Contact & prochaine étape",
        "bullets": [
            "Tester la démo Soundly sur web ou mobile",
            "Créer un salon ou rejoindre un live près de vous",
            "Soutenir un créateur ou explorer Soundly+",
            "Questions & partenariats : contact via l’application",
        ],
        "cta": "Merci — prêts à écouter ensemble ?",
        "image": "05-lives-tab.png",
        "image_caption": "Rejoignez un live ou lancez le vôtre dès maintenant",
    },
]


def img_path(name: str | None) -> Path | None:
    if not name:
        return None
    path = SCREENSHOTS / name
    return path if path.is_file() else None


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, left=Inches(0), top=Inches(0), width=Inches(0.12), height=Inches(7.5)):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_gradient_bar(slide):
    bar = slide.shapes.add_shape(1, Inches(0.6), Inches(1.35), Inches(8.8), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT2
    bar.line.fill.background()


def style_paragraph(p, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.alignment = align
    p.space_after = Pt(8)


def add_title_block(slide, title, top=Inches(0.55)):
    box = slide.shapes.add_textbox(Inches(0.75), top, Inches(8.5), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    style_paragraph(p, size=32, bold=True, color=WHITE)
    add_gradient_bar(slide)
    return box


def add_bullets(slide, items, left=Inches(0.85), top=Inches(1.75), width=Inches(8.3), size=20):
    box = slide.shapes.add_textbox(left, top, width, Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        style_paragraph(p, size=size, color=TEXT)
        p.space_before = Pt(6)
    return box


def add_image_frame(slide, path: Path, left, top, width, caption: str | None = None):
    """Add screenshot with subtle phone-frame border."""
    frame = slide.shapes.add_shape(1, left - Inches(0.06), top - Inches(0.06), width + Inches(0.12), width * 1.95 + Inches(0.12))
    frame.fill.solid()
    frame.fill.fore_color.rgb = BG_ALT
    frame.line.color.rgb = ACCENT
    frame.line.width = Pt(1.5)

    pic_h = width * 1.85
    slide.shapes.add_picture(str(path), left, top, width=width, height=pic_h)

    if caption:
        cap = slide.shapes.add_textbox(left, top + pic_h + Inches(0.08), width + Inches(0.2), Inches(0.55))
        p = cap.text_frame.paragraphs[0]
        p.text = caption
        style_paragraph(p, size=11, color=MUTED, align=PP_ALIGN.CENTER)


def add_slide_image(slide, data, default_left=Inches(0.55), default_top=Inches(1.55), width=Inches(2.65)):
    path = img_path(data.get("image"))
    if not path:
        return False

    layout = data.get("image_layout", "left")
    if layout == "right":
        left = Inches(6.55)
    else:
        left = default_left

    add_image_frame(slide, path, left, default_top, width, data.get("image_caption"))

    secondary = img_path(data.get("image_secondary"))
    if secondary:
        add_image_frame(
            slide,
            secondary,
            left,
            default_top + Inches(3.35),
            width * 0.92,
            None,
        )
    return True


def build_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)

    circle = slide.shapes.add_shape(9, Inches(6.8), Inches(0.4), Inches(2.2), Inches(2.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BG_ALT
    circle.line.color.rgb = ACCENT
    circle.line.width = Pt(2)

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(5.8), Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    p.text = data["title"]
    style_paragraph(p, size=54, bold=True, color=WHITE)

    sub_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.65), Inches(5.8), Inches(0.8))
    p = sub_box.text_frame.paragraphs[0]
    p.text = data["subtitle"]
    style_paragraph(p, size=26, bold=False, color=ACCENT)

    tag_box = slide.shapes.add_textbox(Inches(0.9), Inches(3.55), Inches(5.5), Inches(1.0))
    p = tag_box.text_frame.paragraphs[0]
    p.text = data.get("tagline", "")
    style_paragraph(p, size=17, color=MUTED)

    path = img_path(data.get("image"))
    if path:
        add_image_frame(slide, path, Inches(6.35), Inches(1.35), Inches(2.85), data.get("image_caption"))


def build_highlight_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_block(slide, data["title"])

    has_image = add_slide_image(slide, data, default_top=Inches(1.7), width=Inches(2.5))
    text_left = Inches(0.9)
    text_width = Inches(8.2)
    if has_image and data.get("image_layout") == "right":
        text_width = Inches(5.4)
    elif has_image:
        text_left = Inches(3.55)
        text_width = Inches(5.8)

    box = slide.shapes.add_textbox(text_left, Inches(2.0), text_width, Inches(3.8))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = data["highlight"]
    style_paragraph(p, size=24, color=WHITE, align=PP_ALIGN.LEFT)
    p.line_spacing = 1.35


def build_steps_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_ALT)
    add_accent_bar(slide)
    add_title_block(slide, data["title"])

    has_image = add_slide_image(slide, data, default_left=Inches(6.2), default_top=Inches(1.55), width=Inches(2.55))
    step_width = Inches(5.2) if has_image else Inches(7.5)

    y_start = 1.85
    for i, (num, title, desc) in enumerate(data["steps"]):
        y = y_start + i * 1.15
        badge = slide.shapes.add_shape(9, Inches(0.85), Inches(y), Inches(0.55), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()
        nb = badge.text_frame
        nb.paragraphs[0].text = num
        style_paragraph(nb.paragraphs[0], size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        tbox = slide.shapes.add_textbox(Inches(1.6), Inches(y - 0.05), step_width, Inches(0.45))
        p = tbox.text_frame.paragraphs[0]
        p.text = title
        style_paragraph(p, size=21, bold=True, color=WHITE)

        dbox = slide.shapes.add_textbox(Inches(1.6), Inches(y + 0.35), step_width, Inches(0.4))
        p = dbox.text_frame.paragraphs[0]
        p.text = desc
        style_paragraph(p, size=15, color=MUTED)


def build_bullet_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_accent_bar(slide)
    add_title_block(slide, data["title"])

    has_image = add_slide_image(slide, data)
    bullets_left = Inches(3.55) if has_image else Inches(0.85)
    bullets_width = Inches(5.9) if has_image else Inches(8.3)
    bullets = [f"•  {b}" for b in data.get("bullets", [])]
    add_bullets(slide, bullets, left=bullets_left, top=Inches(1.85), width=bullets_width)

    if data.get("footer"):
        fbox = slide.shapes.add_textbox(Inches(0.85), Inches(6.5), Inches(8.3), Inches(0.6))
        p = fbox.text_frame.paragraphs[0]
        p.text = data["footer"]
        style_paragraph(p, size=16, bold=True, color=ACCENT2)


def build_contact_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_ALT)
    add_accent_bar(slide)
    add_title_block(slide, data["title"])

    has_image = add_slide_image(slide, data, default_left=Inches(6.15), default_top=Inches(1.7), width=Inches(2.6))
    bullets_left = Inches(0.85)
    bullets_width = Inches(5.5) if has_image else Inches(8.3)

    bullets = [f"→  {b}" for b in data.get("bullets", [])]
    add_bullets(slide, bullets, left=bullets_left, top=Inches(1.9), width=bullets_width, size=21)

    cta = slide.shapes.add_textbox(Inches(0.85), Inches(5.8), bullets_width, Inches(0.8))
    p = cta.text_frame.paragraphs[0]
    p.text = data.get("cta", "")
    style_paragraph(p, size=24, bold=True, color=ACCENT, align=PP_ALIGN.LEFT if has_image else PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    missing = []
    for data in SLIDES:
        for key in ("image", "image_secondary"):
            name = data.get(key)
            if name and not img_path(name):
                missing.append(name)

    if missing:
        print(f"Warning: missing screenshots: {', '.join(sorted(set(missing)))}")

    for data in SLIDES:
        t = data.get("type", "bullet")
        if t == "title":
            build_title_slide(prs, data)
        elif "highlight" in data:
            build_highlight_slide(prs, data)
        elif "steps" in data:
            build_steps_slide(prs, data)
        elif t == "contact":
            build_contact_slide(prs, data)
        else:
            build_bullet_slide(prs, data)

    prs.save(str(OUTPUT))
    print(f"Created: {OUTPUT}")
    print(f"Slides: {len(SLIDES)}")
    print(f"Screenshots dir: {SCREENSHOTS}")


if __name__ == "__main__":
    main()
